#!/usr/bin/env python3
"""
MS365 CLI — 通过 COM 自动化操控 Microsoft 365 Word / Excel / PowerPoint / Outlook
对标 wps_cli.py，API 设计兼容，额外支持 Outlook 邮件与日历。

依赖：pip install pywin32 click python-pptx openpyxl
Microsoft 365 需已安装（Office 16+）。
"""

import os
import sys
import json
import time
import click
from pathlib import Path
from typing import Optional
from datetime import datetime, timedelta

# ── COM ProgID ──────────────────────────────────────────────────────────────
MS_APP_IDS = {
    "word":        "Word.Application",
    "excel":       "Excel.Application",
    "powerpoint":  "PowerPoint.Application",
    "outlook":     "Outlook.Application",
}

# PDF 导出常量
WORD_PDF   = 17        # WdExportFormat.wdFormatPDF
EXCEL_PDF  = 0         # xlTypePDF
PPT_PDF    = 32        # PpSaveAsFileType.ppSaveAsPDF
PPT_PNG    = 19        # ppSaveAsPNG


def _get_com_app(app_type: str):
    """获取 MS Office COM 实例（headless，不显示窗口）。

    Excel 使用 DispatchEx 强制新建实例，避免连接到已有的带 UI 的实例导致
    RPC_E_SERVERFAULT（-2146959355）。
    """
    import win32com.client as win32
    prog_id = MS_APP_IDS[app_type]
    try:
        # Excel 单独处理：强制新建实例，再重试 Dispatch
        if app_type == "excel":
            try:
                app = win32.DispatchEx(prog_id)
            except Exception:
                app = win32.Dispatch(prog_id)
        else:
            app = win32.Dispatch(prog_id)

        # PowerPoint: Visible 是只读属性，通过 WithWindow=False 打开来隐藏窗口
        # Outlook: 不设置 Visible，保留用户的运行状态
        if app_type in ("word", "excel"):
            try:
                app.Visible = False
            except Exception:
                pass
            try:
                app.DisplayAlerts = False
            except Exception:
                pass
        return app
    except Exception as e:
        raise RuntimeError(
            f"无法启动 {prog_id}。\n"
            f"请确认 Microsoft 365 已安装。\n"
            f"错误：{e}"
        ) from e


def _find_office_path() -> Optional[str]:
    """从注册表查找 Office 安装路径。"""
    import winreg
    keys_to_try = [
        r"SOFTWARE\Microsoft\Office\16.0\Word\InstallRoot",
        r"SOFTWARE\WOW6432Node\Microsoft\Office\16.0\Word\InstallRoot",
        r"SOFTWARE\Microsoft\Office\15.0\Word\InstallRoot",
    ]
    for hive in [winreg.HKEY_LOCAL_MACHINE, winreg.HKEY_CURRENT_USER]:
        for key_path in keys_to_try:
            try:
                k = winreg.OpenKey(hive, key_path)
                path = winreg.QueryValueEx(k, "Path")[0]
                if os.path.exists(path):
                    return path
            except Exception:
                pass
    return None


def _abs(path: str) -> str:
    return str(Path(path).resolve())


def _out(data: dict, message: str, use_json: bool):
    if use_json:
        click.echo(json.dumps(data, indent=2, ensure_ascii=False))
    else:
        if message:
            click.echo(message)
        for k, v in data.items():
            if not k.startswith("_"):
                click.echo(f"  {k}: {v}")


# ════════════════════════════════════════════════════════════════════════════
# CLI 入口
# ════════════════════════════════════════════════════════════════════════════
@click.group()
@click.option("--json", "use_json", is_flag=True, help="JSON 格式输出")
@click.pass_context
def cli(ctx, use_json):
    """Microsoft 365 CLI — COM 自动化控制 Word / Excel / PowerPoint / Outlook"""
    ctx.ensure_object(dict)
    ctx.obj["json"] = use_json


# ── 检测 ────────────────────────────────────────────────────────────────────
@cli.command()
@click.pass_context
def detect(ctx):
    """检测 Microsoft 365 安装情况和 COM 组件注册状态。"""
    uj = ctx.obj["json"]
    result = {}

    office_path = _find_office_path()
    result["office_path"] = office_path or "未找到（注册表查询失败）"

    import win32com.client as win32
    import winreg

    # Outlook：只检查注册表，不实际启动（启动耗时且可能弹窗）
    try:
        k = winreg.OpenKey(winreg.HKEY_LOCAL_MACHINE,
                           r"SOFTWARE\Microsoft\Windows\CurrentVersion\App Paths\OUTLOOK.EXE")
        outlook_exe = winreg.QueryValueEx(k, "")[0]
        result["com_outlook"] = f"✓ Outlook.Application  exe={outlook_exe}"
    except Exception:
        result["com_outlook"] = "✗ Outlook 未找到（注册表查询失败）"

    for app_type, prog_id in MS_APP_IDS.items():
        if app_type == "outlook":
            continue  # 已在上方处理
        try:
            app = (win32.DispatchEx(prog_id) if app_type == "excel"
                   else win32.Dispatch(prog_id))
            if app_type in ("word", "excel"):
                app.Visible = False
                app.DisplayAlerts = False
            ver = getattr(app, "Version", "?")
            result[f"com_{app_type}"] = f"✓ {prog_id}  Office {ver}"
            try:
                app.Quit()
            except Exception:
                pass
        except Exception as e:
            result[f"com_{app_type}"] = f"✗ {e}"

    _out(result, "Microsoft 365 检测结果", uj)


# ════════════════════════════════════════════════════════════════════════════
# Word
# ════════════════════════════════════════════════════════════════════════════
@cli.group()
def word():
    """Microsoft Word 文档操作。"""


@word.command("new")
@click.option("-o", "--output", required=True, help="保存路径 (.docx)")
@click.option("-t", "--title", default="新建文档", help="文档标题")
@click.option("--body", default="", help="初始正文内容")
@click.pass_context
def word_new(ctx, output, title, body):
    """新建 Word 文档。"""
    uj = ctx.obj["json"]
    output = _abs(output)
    try:
        app = _get_com_app("word")
        doc = app.Documents.Add()
        sel = app.Selection
        # 用字体大小模拟标题样式（避免样式名语言差异）
        sel.Font.Size = 22
        sel.Font.Bold = True
        sel.ParagraphFormat.Alignment = 1  # 居中
        sel.TypeText(title)
        sel.TypeParagraph()
        if body:
            sel.Font.Size = 12
            sel.Font.Bold = False
            sel.ParagraphFormat.Alignment = 0
            sel.TypeText(body)
        doc.SaveAs2(output)
        doc.Close(False)
        app.Quit()
        size = os.path.getsize(output)
        _out({"path": output, "title": title, "size_bytes": size},
             f"✓ 已创建: {output}", uj)
    except Exception as e:
        click.echo(f"错误: {e}", err=True); sys.exit(1)


@word.command("add-text")
@click.argument("docpath")
@click.option("-t", "--text", required=True, help="要追加的文字")
@click.option("--heading", type=int, default=0, help="标题级别(1-9)，0=正文")
@click.pass_context
def word_add_text(ctx, docpath, text, heading):
    """向现有 Word 文档追加文字。"""
    uj = ctx.obj["json"]
    docpath = _abs(docpath)
    try:
        app = _get_com_app("word")
        doc = app.Documents.Open(docpath)
        sel = app.Selection
        sel.EndKey(6)  # wdStory = 6，移到文档末尾
        # WdBuiltinStyle: wdStyleHeading1=-2 ... wdStyleHeading9=-10, wdStyleNormal=-1
        if heading > 0:
            heading_const = -(1 + heading)  # Heading1=-2, Heading2=-3, ...
            try:
                sel.Style = doc.Styles(heading_const)
            except Exception:
                # 回退：用字体大小区分级别
                sel.Font.Size = max(22 - (heading - 1) * 2, 10)
                sel.Font.Bold = True
        else:
            try:
                sel.Style = doc.Styles(-1)  # wdStyleNormal
            except Exception:
                sel.Font.Size = 12
                sel.Font.Bold = False
        sel.TypeText(text)
        sel.TypeParagraph()
        doc.Save()
        doc.Close(False)
        app.Quit()
        _out({"path": docpath, "added": text[:80], "heading": heading},
             f"✓ 已追加文字", uj)
    except Exception as e:
        click.echo(f"错误: {e}", err=True); sys.exit(1)


@word.command("read")
@click.argument("docpath")
@click.option("--max-chars", default=5000, help="最多读取字符数（默认 5000）")
@click.pass_context
def word_read(ctx, docpath, max_chars):
    """读取 Word 文档文字内容。"""
    uj = ctx.obj["json"]
    docpath = _abs(docpath)
    try:
        app = _get_com_app("word")
        doc = app.Documents.Open(docpath, ReadOnly=True)
        text = doc.Content.Text[:max_chars]
        word_count = doc.Words.Count
        para_count = doc.Paragraphs.Count
        doc.Close(False)
        app.Quit()
        _out({"path": docpath, "words": word_count, "paragraphs": para_count,
              "text_preview": text},
             f"✓ 文档内容（前 {max_chars} 字符）:\n{text}", uj)
    except Exception as e:
        click.echo(f"错误: {e}", err=True); sys.exit(1)


@word.command("find-replace")
@click.argument("docpath")
@click.option("--find", "find_text", required=True, help="要查找的文字")
@click.option("--replace", "replace_text", required=True, help="替换为")
@click.option("--all", "replace_all", is_flag=True, default=True, help="替换所有（默认是）")
@click.pass_context
def word_find_replace(ctx, docpath, find_text, replace_text, replace_all):
    """文档内查找并替换文字。"""
    uj = ctx.obj["json"]
    docpath = _abs(docpath)
    try:
        app = _get_com_app("word")
        doc = app.Documents.Open(docpath)
        # 用文档全范围（Content）做查找替换，避免 Selection 范围问题
        rng = doc.Content
        rng.Find.ClearFormatting()
        rng.Find.Replacement.ClearFormatting()
        # Execute(FindText, MatchCase, MatchWholeWord, MatchWildcards,
        #         MatchSoundsLike, MatchAllWordForms, Forward, Wrap,
        #         Format, ReplaceWith, Replace)
        rng.Find.Execute(
            find_text, False, False, False, False, False,
            True, 1, False,
            replace_text,
            2 if replace_all else 1
        )
        doc.Save()
        doc.Close(False)
        app.Quit()
        _out({"path": docpath, "find": find_text, "replace": replace_text,
              "replace_all": replace_all},
             f"✓ 替换完成: '{find_text}' → '{replace_text}'", uj)
    except Exception as e:
        click.echo(f"错误: {e}", err=True); sys.exit(1)


@word.command("to-pdf")
@click.argument("docpath")
@click.option("-o", "--output", default=None, help="PDF 输出路径")
@click.pass_context
def word_to_pdf(ctx, docpath, output):
    """将 Word 文档导出为 PDF。"""
    uj = ctx.obj["json"]
    docpath = _abs(docpath)
    if not output:
        output = str(Path(docpath).with_suffix(".pdf"))
    output = _abs(output)
    try:
        app = _get_com_app("word")
        doc = app.Documents.Open(docpath)
        doc.ExportAsFixedFormat(output, WORD_PDF)
        doc.Close(False)
        app.Quit()
        size = os.path.getsize(output)
        with open(output, "rb") as f:
            valid = f.read(4) == b"%PDF"
        _out({"input": docpath, "output": output,
              "size_bytes": size, "valid_pdf": valid},
             f"✓ PDF: {output} ({size:,} bytes)", uj)
    except Exception as e:
        click.echo(f"错误: {e}", err=True); sys.exit(1)


# ════════════════════════════════════════════════════════════════════════════
# Excel
# ════════════════════════════════════════════════════════════════════════════
@cli.group()
def excel():
    """Microsoft Excel 表格操作。"""


@excel.command("new")
@click.option("-o", "--output", required=True, help="保存路径 (.xlsx)")
@click.option("--sheet", default="Sheet1", help="首个工作表名称")
@click.pass_context
def excel_new(ctx, output, sheet):
    """新建 Excel 工作簿。"""
    uj = ctx.obj["json"]
    output = _abs(output)
    try:
        app = _get_com_app("excel")
        wb = app.Workbooks.Add()
        wb.Worksheets(1).Name = sheet
        # 删除多余工作表
        while wb.Worksheets.Count > 1:
            wb.Worksheets(wb.Worksheets.Count).Delete()
        wb.SaveAs(output)
        wb.Close(False)
        app.Quit()
        _out({"path": output, "sheet": sheet}, f"✓ 已创建: {output}", uj)
    except Exception as e:
        click.echo(f"错误: {e}", err=True); sys.exit(1)


@excel.command("write-cell")
@click.argument("xlspath")
@click.option("--cell", required=True, help="单元格地址，如 A1 或 B3")
@click.option("--value", required=True, help="写入值")
@click.option("--sheet", default="1", type=str, help="工作表序号或名称（默认 1）")
@click.pass_context
def excel_write_cell(ctx, xlspath, cell, value, sheet):
    """写入指定单元格。"""
    uj = ctx.obj["json"]
    xlspath = _abs(xlspath)
    try:
        app = _get_com_app("excel")
        wb = app.Workbooks.Open(xlspath)
        ws = wb.Worksheets(int(sheet) if str(sheet).isdigit() else sheet)
        # 尝试转换数字
        try:
            ws.Range(cell).Value = float(value) if "." in value else int(value)
        except ValueError:
            ws.Range(cell).Value = value
        wb.Save()
        wb.Close(False)
        app.Quit()
        _out({"file": xlspath, "cell": cell, "value": value, "sheet": sheet},
             f"✓ [{sheet}]{cell} = {value}", uj)
    except Exception as e:
        click.echo(f"错误: {e}", err=True); sys.exit(1)


@excel.command("read-range")
@click.argument("xlspath")
@click.option("--range", "cell_range", required=True, help="范围，如 A1:C5 或 A1")
@click.option("--sheet", default="1", type=str, help="工作表序号或名称（默认 1）")
@click.pass_context
def excel_read_range(ctx, xlspath, cell_range, sheet):
    """读取单元格范围数据。"""
    uj = ctx.obj["json"]
    xlspath = _abs(xlspath)
    try:
        app = _get_com_app("excel")
        wb = app.Workbooks.Open(xlspath, ReadOnly=True)
        ws = wb.Worksheets(int(sheet) if str(sheet).isdigit() else sheet)
        rng = ws.Range(cell_range)
        # 单格 vs 区域
        if rng.Rows.Count == 1 and rng.Columns.Count == 1:
            data = [[rng.Value]]
        else:
            raw = rng.Value
            data = [list(row) for row in raw] if raw else []
        wb.Close(False)
        app.Quit()
        _out({"file": xlspath, "range": cell_range, "sheet": sheet,
              "rows": len(data), "data": data},
             f"✓ 读取 {cell_range}（{len(data)} 行）", uj)
    except Exception as e:
        click.echo(f"错误: {e}", err=True); sys.exit(1)


@excel.command("run-macro")
@click.argument("xlspath")
@click.option("--macro", required=True, help="宏名称（如 Module1.MyMacro）")
@click.option("--args", default="", help="宏参数（逗号分隔）")
@click.pass_context
def excel_run_macro(ctx, xlspath, macro, args):
    """执行工作簿中的 VBA 宏。"""
    uj = ctx.obj["json"]
    xlspath = _abs(xlspath)
    try:
        app = _get_com_app("excel")
        wb = app.Workbooks.Open(xlspath)
        macro_args = [a.strip() for a in args.split(",") if a.strip()] if args else []
        if macro_args:
            app.Run(macro, *macro_args)
        else:
            app.Run(macro)
        wb.Save()
        wb.Close(False)
        app.Quit()
        _out({"file": xlspath, "macro": macro, "args": macro_args},
             f"✓ 宏已执行: {macro}", uj)
    except Exception as e:
        click.echo(f"错误: {e}", err=True); sys.exit(1)


@excel.command("to-pdf")
@click.argument("xlspath")
@click.option("-o", "--output", default=None, help="PDF 输出路径")
@click.option("--sheet", default=None, type=str, help="仅导出指定工作表（默认全部）")
@click.pass_context
def excel_to_pdf(ctx, xlspath, output, sheet):
    """将 Excel 工作簿（或指定工作表）导出为 PDF。"""
    uj = ctx.obj["json"]
    xlspath = _abs(xlspath)
    if not output:
        output = str(Path(xlspath).with_suffix(".pdf"))
    output = _abs(output)
    try:
        app = _get_com_app("excel")
        wb = app.Workbooks.Open(xlspath)
        if sheet:
            ws = wb.Worksheets(int(sheet) if str(sheet).isdigit() else sheet)
            ws.ExportAsFixedFormat(EXCEL_PDF, output)
        else:
            wb.ExportAsFixedFormat(EXCEL_PDF, output)
        wb.Close(False)
        app.Quit()
        size = os.path.getsize(output)
        _out({"input": xlspath, "output": output, "size_bytes": size},
             f"✓ PDF: {output} ({size:,} bytes)", uj)
    except Exception as e:
        click.echo(f"错误: {e}", err=True); sys.exit(1)


# ════════════════════════════════════════════════════════════════════════════
# PowerPoint
# ════════════════════════════════════════════════════════════════════════════
@cli.group()
def powerpoint():
    """Microsoft PowerPoint 演示文稿操作。"""


@powerpoint.command("new")
@click.option("-o", "--output", required=True, help="保存路径 (.pptx)")
@click.option("--title", default="新建演示", help="首页标题")
@click.option("--subtitle", default="", help="首页副标题")
@click.pass_context
def powerpoint_new(ctx, output, title, subtitle):
    """新建 PowerPoint 演示文稿。"""
    uj = ctx.obj["json"]
    output = _abs(output)
    try:
        from pptx import Presentation
        from pptx.util import Pt
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        if subtitle and len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle
        prs.save(output)
        size = os.path.getsize(output)
        _out({"path": output, "title": title, "slides": 1, "size_bytes": size},
             f"✓ 已创建: {output}", uj)
    except Exception as e:
        click.echo(f"错误: {e}", err=True); sys.exit(1)


@powerpoint.command("add-slide")
@click.argument("pptpath")
@click.option("-t", "--title", required=True, help="幻灯片标题")
@click.option("--body", default="", help="正文内容")
@click.option("--layout", default=1, type=int, help="布局序号（0=标题页，1=标题+内容，6=空白）")
@click.pass_context
def powerpoint_add_slide(ctx, pptpath, title, body, layout):
    """向 PPT 追加幻灯片。"""
    uj = ctx.obj["json"]
    pptpath = _abs(pptpath)
    try:
        from pptx import Presentation
        prs = Presentation(pptpath)
        lyt = prs.slide_layouts[layout]
        slide = prs.slides.add_slide(lyt)
        if slide.shapes.title:
            slide.shapes.title.text = title
        if body and len(slide.placeholders) > 1:
            slide.placeholders[1].text = body
        prs.save(pptpath)
        _out({"path": pptpath, "title": title,
              "total_slides": len(prs.slides), "layout": layout},
             f"✓ 已添加幻灯片（共 {len(prs.slides)} 张）", uj)
    except Exception as e:
        click.echo(f"错误: {e}", err=True); sys.exit(1)


@powerpoint.command("to-pdf")
@click.argument("pptpath")
@click.option("-o", "--output", default=None, help="PDF 输出路径")
@click.pass_context
def powerpoint_to_pdf(ctx, pptpath, output):
    """将 PPT 通过 Office COM 导出为 PDF。"""
    uj = ctx.obj["json"]
    pptpath = _abs(pptpath)
    if not output:
        output = str(Path(pptpath).with_suffix(".pdf"))
    output = _abs(output)
    try:
        app = _get_com_app("powerpoint")
        prs = app.Presentations.Open(pptpath, WithWindow=False)
        prs.SaveAs(output, PPT_PDF)
        prs.Close()
        app.Quit()
        size = os.path.getsize(output)
        with open(output, "rb") as f:
            valid = f.read(4) == b"%PDF"
        _out({"input": pptpath, "output": output,
              "size_bytes": size, "valid_pdf": valid},
             f"✓ PDF: {output} ({size:,} bytes)", uj)
    except Exception as e:
        click.echo(f"错误: {e}", err=True); sys.exit(1)


@powerpoint.command("to-images")
@click.argument("pptpath")
@click.option("-o", "--outdir", required=True, help="图片输出目录")
@click.option("--format", "fmt", default="png", type=click.Choice(["png", "jpg"]),
              help="图片格式（默认 png）")
@click.option("--width", default=1920, type=int, help="图片宽度（像素，默认 1920）")
@click.option("--height", default=1080, type=int, help="图片高度（像素，默认 1080）")
@click.pass_context
def powerpoint_to_images(ctx, pptpath, outdir, fmt, width, height):
    """将 PPT 每张幻灯片导出为图片。"""
    uj = ctx.obj["json"]
    pptpath = _abs(pptpath)
    outdir = _abs(outdir)
    Path(outdir).mkdir(parents=True, exist_ok=True)
    try:
        app = _get_com_app("powerpoint")
        prs = app.Presentations.Open(pptpath, WithWindow=False)
        slides_count = prs.Slides.Count
        saved = []
        for i in range(1, slides_count + 1):
            img_path = os.path.join(outdir, f"slide_{i:03d}.{fmt}")
            prs.Slides(i).Export(img_path, fmt.upper(), width, height)
            saved.append(img_path)
        prs.Close()
        app.Quit()
        _out({"input": pptpath, "outdir": outdir, "slides": slides_count,
              "format": fmt, "files": saved},
             f"✓ 已导出 {slides_count} 张幻灯片 → {outdir}", uj)
    except Exception as e:
        click.echo(f"错误: {e}", err=True); sys.exit(1)


# ════════════════════════════════════════════════════════════════════════════
# Outlook
# ════════════════════════════════════════════════════════════════════════════
@cli.group()
def outlook():
    """Microsoft Outlook 邮件与日历操作。"""


def _get_outlook_namespace():
    """获取 Outlook MAPI namespace。"""
    app = _get_com_app("outlook")
    ns = app.GetNamespace("MAPI")
    return app, ns


@outlook.command("send")
@click.option("--to", "to_addr", required=True, help="收件人邮件地址（多个用分号分隔）")
@click.option("--subject", required=True, help="邮件主题")
@click.option("--body", default="", help="纯文本正文")
@click.option("--html", "html_body", default="", help="HTML 正文（优先级高于 --body）")
@click.option("--attach", multiple=True, help="附件路径（可多次指定）")
@click.option("--cc", default="", help="抄送（多个用分号分隔）")
@click.option("--bcc", default="", help="密送（多个用分号分隔）")
@click.option("--draft", is_flag=True, help="只保存为草稿，不发送")
@click.pass_context
def outlook_send(ctx, to_addr, subject, body, html_body, attach, cc, bcc, draft):
    """发送邮件（或保存草稿）。"""
    uj = ctx.obj["json"]
    try:
        app = _get_com_app("outlook")
        mail = app.CreateItem(0)  # 0 = olMailItem
        mail.To = to_addr
        mail.Subject = subject
        if cc:
            mail.CC = cc
        if bcc:
            mail.BCC = bcc
        if html_body:
            mail.HTMLBody = html_body
        else:
            mail.Body = body
        for att in attach:
            att_abs = _abs(att)
            if not os.path.exists(att_abs):
                click.echo(f"警告: 附件不存在 {att_abs}", err=True)
            else:
                mail.Attachments.Add(att_abs)
        if draft:
            mail.Save()
            action = "草稿已保存"
        else:
            mail.Send()
            action = "邮件已发送"
        _out({"to": to_addr, "subject": subject, "action": action,
              "attachments": len(attach)},
             f"✓ {action}: {subject} → {to_addr}", uj)
    except Exception as e:
        click.echo(f"错误: {e}", err=True); sys.exit(1)


@outlook.command("list-inbox")
@click.option("--count", default=10, type=int, help="显示最新邮件数量（默认 10）")
@click.option("--unread-only", is_flag=True, help="只显示未读邮件")
@click.pass_context
def outlook_list_inbox(ctx, count, unread_only):
    """列出收件箱邮件。"""
    uj = ctx.obj["json"]
    try:
        app, ns = _get_outlook_namespace()
        inbox = ns.GetDefaultFolder(6)  # 6 = olFolderInbox
        items = inbox.Items
        items.Sort("[ReceivedTime]", True)  # 降序（最新在前）
        mails = []
        collected = 0
        for item in items:
            if collected >= count:
                break
            try:
                unread = not item.UnRead if hasattr(item, "UnRead") else False
                # UnRead=True 表示未读
                is_unread = getattr(item, "UnRead", False)
                if unread_only and not is_unread:
                    continue
                mails.append({
                    "index": collected + 1,
                    "subject": getattr(item, "Subject", ""),
                    "sender": getattr(item, "SenderEmailAddress", ""),
                    "received": str(getattr(item, "ReceivedTime", "")),
                    "unread": is_unread,
                    "has_attachments": getattr(item, "Attachments", None) and
                                       item.Attachments.Count > 0,
                })
                collected += 1
            except Exception:
                continue
        _out({"inbox": mails, "shown": len(mails)},
             f"✓ 收件箱（显示 {len(mails)} 封）", uj)
    except Exception as e:
        click.echo(f"错误: {e}", err=True); sys.exit(1)


@outlook.group()
def calendar():
    """日历操作。"""


@calendar.command("list")
@click.option("--days", default=7, type=int, help="查询未来天数（默认 7）")
@click.option("--past", default=0, type=int, help="同时查询过去天数（默认 0）")
@click.pass_context
def calendar_list(ctx, days, past):
    """列出日历事件。"""
    uj = ctx.obj["json"]
    try:
        app, ns = _get_outlook_namespace()
        cal = ns.GetDefaultFolder(9)  # 9 = olFolderCalendar
        items = cal.Items
        items.IncludeRecurrences = True
        items.Sort("[Start]")

        now = datetime.now()
        start_dt = now - timedelta(days=past)
        end_dt = now + timedelta(days=days)

        # Outlook 日期过滤字符串
        fmt = "%m/%d/%Y %H:%M %p"
        restriction = (
            f"[Start] >= '{start_dt.strftime(fmt)}' "
            f"AND [Start] <= '{end_dt.strftime(fmt)}'"
        )
        try:
            filtered = items.Restrict(restriction)
        except Exception:
            filtered = items  # 过滤失败则显示全部

        events = []
        for item in filtered:
            try:
                events.append({
                    "subject": getattr(item, "Subject", ""),
                    "start": str(getattr(item, "Start", "")),
                    "end": str(getattr(item, "End", "")),
                    "location": getattr(item, "Location", ""),
                    "organizer": getattr(item, "Organizer", ""),
                    "all_day": getattr(item, "AllDayEvent", False),
                })
            except Exception:
                continue

        _out({"events": events, "total": len(events),
              "from": str(start_dt.date()), "to": str(end_dt.date())},
             f"✓ 日历事件（{len(events)} 个，未来 {days} 天）", uj)
    except Exception as e:
        click.echo(f"错误: {e}", err=True); sys.exit(1)


@calendar.command("add")
@click.option("--subject", required=True, help="事件主题")
@click.option("--start", required=True, help="开始时间，如 '2026-03-15 14:00'")
@click.option("--end", required=True, help="结束时间，如 '2026-03-15 15:00'")
@click.option("--location", default="", help="地点")
@click.option("--body", default="", help="备注内容")
@click.option("--all-day", is_flag=True, help="全天事件")
@click.option("--reminder", default=15, type=int, help="提醒分钟数（默认 15，0=不提醒）")
@click.pass_context
def calendar_add(ctx, subject, start, end, location, body, all_day, reminder):
    """新建日历事件。"""
    uj = ctx.obj["json"]
    try:
        # 解析时间
        for fmt in ["%Y-%m-%d %H:%M", "%Y-%m-%d %H:%M:%S", "%Y-%m-%d"]:
            try:
                start_dt = datetime.strptime(start, fmt)
                end_dt = datetime.strptime(end, fmt)
                break
            except ValueError:
                continue
        else:
            raise ValueError(f"无法解析时间格式: {start}")

        app = _get_com_app("outlook")
        appt = app.CreateItem(1)  # 1 = olAppointmentItem
        appt.Subject = subject
        appt.Start = start_dt.strftime("%m/%d/%Y %H:%M")
        appt.End = end_dt.strftime("%m/%d/%Y %H:%M")
        if location:
            appt.Location = location
        if body:
            appt.Body = body
        appt.AllDayEvent = all_day
        if reminder > 0:
            appt.ReminderSet = True
            appt.ReminderMinutesBeforeStart = reminder
        else:
            appt.ReminderSet = False
        appt.Save()

        _out({"subject": subject, "start": str(start_dt),
              "end": str(end_dt), "location": location},
             f"✓ 已添加日历事件: {subject} @ {start}", uj)
    except Exception as e:
        click.echo(f"错误: {e}", err=True); sys.exit(1)


# ════════════════════════════════════════════════════════════════════════════
# 通用转换
# ════════════════════════════════════════════════════════════════════════════
@cli.command("convert")
@click.argument("input_path")
@click.argument("output_path")
@click.pass_context
def convert(ctx, input_path, output_path):
    """智能格式转换（自动识别文档类型）。

    示例：ms365-cli convert report.docx report.pdf
    """
    ext = Path(input_path).suffix.lower()
    if ext in (".docx", ".doc", ".odt", ".rtf", ".txt"):
        ctx.invoke(word_to_pdf, docpath=input_path, output=output_path)
    elif ext in (".xlsx", ".xls", ".csv"):
        ctx.invoke(excel_to_pdf, xlspath=input_path, output=output_path)
    elif ext in (".pptx", ".ppt"):
        ctx.invoke(powerpoint_to_pdf, pptpath=input_path, output=output_path)
    else:
        click.echo(f"不支持的格式: {ext}", err=True); sys.exit(1)


# ── 版本信息 ──────────────────────────────────────────────────────────────
@cli.command("version")
@click.pass_context
def version(ctx):
    """显示 Microsoft Office 版本信息。"""
    uj = ctx.obj["json"]
    try:
        app = _get_com_app("word")
        ver = app.Version
        build = getattr(app, "Build", "N/A")
        app.Quit()
        # 版本号映射
        ver_map = {
            "16.0": "Office 2016/2019/2021/365",
            "15.0": "Office 2013",
            "14.0": "Office 2010",
        }
        ver_name = ver_map.get(ver, f"Office ({ver})")
        _out({"version": ver, "build": build, "product": ver_name,
              "com_id": "Word.Application"},
             f"Microsoft {ver_name}  (v{ver} build {build})", uj)
    except Exception as e:
        click.echo(f"错误: {e}", err=True)


if __name__ == "__main__":
    cli(obj={})
