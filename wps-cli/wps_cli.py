#!/usr/bin/env python3
"""
WPS Office CLI — 通过 COM 自动化操控 WPS Writer / 表格 / 演示
对标 cli-anything-libreoffice，支持相同的核心功能。

依赖：pip install pywin32 click
WPS 需已安装并注册 COM 组件。
"""

import os
import sys
import json
import time
import click
from pathlib import Path
from typing import Optional

# ── COM 组件 ID ─────────────────────────────────────────────────
WPS_APP_IDS = {
    "writer":  ["KWPS.Application",  "WPS.Application"],
    "calc":    ["KET.Application",   "ET.Application"],
    # KWPP 注册在 HKCU，需用 CLSID 兜底
    "impress": ["KWPP.Application", "KWpp.Application",
                "{44720441-94BF-4940-926D-4F38FECF2A48}", "WPP.Application"],
}


def _clsid_from_progid(prog_id: str) -> Optional[str]:
    """从注册表（HKCU 优先，再 HKLM）查找 ProgID 对应的 CLSID。"""
    import winreg
    for hive in [winreg.HKEY_CURRENT_USER, winreg.HKEY_LOCAL_MACHINE]:
        try:
            k = winreg.OpenKey(hive, f"SOFTWARE\\Classes\\{prog_id}\\CLSID")
            return winreg.QueryValueEx(k, "")[0]
        except Exception:
            pass
    return None


def _get_com_app(doc_type: str):
    """获取 WPS COM 应用实例，支持 HKCU/HKLM 双注册方式。"""
    import win32com.client as win32
    import pythoncom, pywintypes
    ids = WPS_APP_IDS.get(doc_type, WPS_APP_IDS["writer"])
    last_err = None
    for prog_id in ids:
        try:
            # 先尝试标准 Dispatch（HKLM 注册的 ProgID 直接可用）
            app = win32.Dispatch(prog_id)
            app.Visible = False
            return app
        except Exception:
            pass
        # 标准 Dispatch 失败时，从注册表找 CLSID 再 CoCreateInstance（兼容 HKCU）
        try:
            clsid_str = _clsid_from_progid(prog_id)
            if not clsid_str:
                continue
            clsid = pywintypes.IID(clsid_str)
            obj = pythoncom.CoCreateInstance(
                clsid, None, pythoncom.CLSCTX_ALL, pythoncom.IID_IDispatch
            )
            app = win32.Dispatch(obj)
            app.Visible = False
            return app
        except Exception as e:
            last_err = e
    raise RuntimeError(
        f"无法启动 WPS COM（{doc_type}）。\n"
        f"请确认 WPS 已安装且注册了 COM 组件。\n"
        f"最后错误：{last_err}"
    )


def _find_wps_exe():
    """自动查找 WPS 可执行文件路径。"""
    candidates = [
        r"C:\Program Files (x86)\Kingsoft\WPS Office\office6\wps.exe",
        r"C:\Program Files\Kingsoft\WPS Office\office6\wps.exe",
    ]
    import winreg
    try:
        key = winreg.OpenKey(winreg.HKEY_LOCAL_MACHINE,
                             r"SOFTWARE\WOW6432Node\Kingsoft\WPS\wps")
        path = winreg.QueryValueEx(key, "InstallPath")[0]
        candidates.insert(0, os.path.join(path, "office6", "wps.exe"))
    except Exception:
        pass
    for c in candidates:
        if os.path.exists(c):
            return c
    # 动态搜索
    import glob
    hits = glob.glob(r"C:\**\Kingsoft\**\wps.exe", recursive=True)
    if hits:
        return hits[0]
    return None


def _abs(path: str) -> str:
    return str(Path(path).resolve())


def _out(data: dict, message: str, use_json: bool):
    if use_json:
        click.echo(json.dumps(data, indent=2, ensure_ascii=False))
    else:
        if message:
            click.echo(message)
        for k, v in data.items():
            click.echo(f"  {k}: {v}")


# ════════════════════════════════════════════════════════════════
# CLI 入口
# ════════════════════════════════════════════════════════════════
@click.group()
@click.option("--json", "use_json", is_flag=True, help="JSON 格式输出")
@click.pass_context
def cli(ctx, use_json):
    """WPS Office CLI — COM 自动化控制 WPS Writer / 表格 / 演示"""
    ctx.ensure_object(dict)
    ctx.obj["json"] = use_json


# ── 检测 ────────────────────────────────────────────────────────
@cli.command()
@click.pass_context
def detect(ctx):
    """检测 WPS 安装情况和 COM 组件注册状态。"""
    uj = ctx.obj["json"]
    result = {}

    exe = _find_wps_exe()
    result["wps_exe"] = exe or "未找到"
    result["exe_exists"] = bool(exe)

    import win32com.client as win32
    for dtype, ids in WPS_APP_IDS.items():
        for pid in ids:
            try:
                app = win32.Dispatch(pid)
                app.Quit()
                result[f"com_{dtype}"] = f"✓ {pid}"
                break
            except Exception as e:
                result[f"com_{dtype}"] = f"✗ {e}"

    _out(result, "WPS 检测结果", uj)


# ── Writer 文档 ─────────────────────────────────────────────────
@cli.group()
def writer():
    """WPS Writer 文档操作。"""


@writer.command("new")
@click.option("-o", "--output", required=True, help="保存路径 (.docx)")
@click.option("-t", "--title", default="新建文档", help="文档标题")
@click.pass_context
def writer_new(ctx, output, title):
    """新建 Writer 文档。"""
    uj = ctx.obj["json"]
    output = _abs(output)
    try:
        app = _get_com_app("writer")
        doc = app.Documents.Add()
        sel = app.Selection
        sel.ParagraphFormat.Alignment = 1  # 居中
        sel.Font.Size = 22
        sel.Font.Bold = True
        sel.TypeText(title)
        sel.TypeParagraph()
        sel.Font.Size = 12
        sel.Font.Bold = False
        sel.ParagraphFormat.Alignment = 0
        doc.SaveAs(output)
        doc.Close(False)
        app.Quit()
        size = os.path.getsize(output)
        _out({"path": output, "title": title, "size_bytes": size}, f"✓ 已创建: {output}", uj)
    except Exception as e:
        click.echo(f"错误: {e}", err=True); sys.exit(1)


@writer.command("add-text")
@click.argument("docpath")
@click.option("-t", "--text", required=True, help="要追加的文字")
@click.option("--heading", type=int, default=0, help="标题级别(1-6)，0=正文")
@click.pass_context
def writer_add_text(ctx, docpath, text, heading):
    """向现有文档追加文字。"""
    uj = ctx.obj["json"]
    docpath = _abs(docpath)
    try:
        app = _get_com_app("writer")
        doc = app.Documents.Open(docpath)
        sel = app.Selection
        sel.EndKey(6)  # 移到末尾
        if heading > 0:
            sel.Style = f"标题 {heading}"
        sel.TypeText(text)
        sel.TypeParagraph()
        sel.Style = "正文"
        doc.Save()
        doc.Close(False)
        app.Quit()
        _out({"path": docpath, "added": text, "heading": heading}, f"✓ 已追加文字", uj)
    except Exception as e:
        click.echo(f"错误: {e}", err=True); sys.exit(1)


@writer.command("to-pdf")
@click.argument("docpath")
@click.option("-o", "--output", default=None, help="PDF 输出路径")
@click.pass_context
def writer_to_pdf(ctx, docpath, output):
    """将 Writer 文档导出为 PDF（调用 WPS 真实渲染）。"""
    uj = ctx.obj["json"]
    docpath = _abs(docpath)
    if not output:
        output = str(Path(docpath).with_suffix(".pdf"))
    output = _abs(output)
    try:
        app = _get_com_app("writer")
        doc = app.Documents.Open(docpath)
        # WPS ExportAsFixedFormat: 17=PDF
        doc.ExportAsFixedFormat(output, 17)
        doc.Close(False)
        app.Quit()
        size = os.path.getsize(output)
        # 验证 PDF 魔术字节
        with open(output, "rb") as f:
            magic = f.read(4)
        valid = magic == b"%PDF"
        _out({"input": docpath, "output": output,
              "size_bytes": size, "valid_pdf": valid},
             f"✓ PDF 已生成: {output} ({size:,} bytes)", uj)
    except Exception as e:
        click.echo(f"错误: {e}", err=True); sys.exit(1)


# ── 表格 ────────────────────────────────────────────────────────
@cli.group()
def calc():
    """WPS 表格（ET）操作。"""


@calc.command("new")
@click.option("-o", "--output", required=True, help="保存路径 (.xlsx)")
@click.option("--sheet", default="Sheet1", help="工作表名")
@click.pass_context
def calc_new(ctx, output, sheet):
    """新建表格文件。"""
    uj = ctx.obj["json"]
    output = _abs(output)
    try:
        app = _get_com_app("calc")
        wb = app.Workbooks.Add()
        ws = wb.Worksheets(1)
        ws.Name = sheet
        wb.SaveAs(output)
        wb.Close(False)
        app.Quit()
        _out({"path": output, "sheet": sheet}, f"✓ 已创建: {output}", uj)
    except Exception as e:
        click.echo(f"错误: {e}", err=True); sys.exit(1)


@calc.command("write-cell")
@click.argument("xlspath")
@click.option("--cell", required=True, help="单元格地址，如 A1")
@click.option("--value", required=True, help="写入值")
@click.option("--sheet", default=1, help="工作表序号")
@click.pass_context
def calc_write_cell(ctx, xlspath, cell, value, sheet):
    """写入单元格数据。"""
    uj = ctx.obj["json"]
    xlspath = _abs(xlspath)
    try:
        app = _get_com_app("calc")
        wb = app.Workbooks.Open(xlspath)
        ws = wb.Worksheets(sheet)
        ws.Range(cell).Value = value
        wb.Save()
        wb.Close(False)
        app.Quit()
        _out({"file": xlspath, "cell": cell, "value": value}, f"✓ {cell} = {value}", uj)
    except Exception as e:
        click.echo(f"错误: {e}", err=True); sys.exit(1)


@calc.command("to-pdf")
@click.argument("xlspath")
@click.option("-o", "--output", default=None)
@click.pass_context
def calc_to_pdf(ctx, xlspath, output):
    """将表格导出为 PDF。"""
    uj = ctx.obj["json"]
    xlspath = _abs(xlspath)
    if not output:
        output = str(Path(xlspath).with_suffix(".pdf"))
    output = _abs(output)
    try:
        app = _get_com_app("calc")
        wb = app.Workbooks.Open(xlspath)
        wb.ExportAsFixedFormat(0, output)  # 0=xlTypePDF
        wb.Close(False)
        app.Quit()
        size = os.path.getsize(output)
        _out({"input": xlspath, "output": output, "size_bytes": size},
             f"✓ PDF: {output}", uj)
    except Exception as e:
        click.echo(f"错误: {e}", err=True); sys.exit(1)


# ── 演示文稿 ────────────────────────────────────────────────────
@cli.group()
def impress():
    """WPS 演示（Presentation）操作。"""


@impress.command("new")
@click.option("-o", "--output", required=True, help="保存路径 (.pptx)")
@click.option("--title", default="新建演示", help="首页标题")
@click.option("--subtitle", default="", help="副标题")
@click.pass_context
def impress_new(ctx, output, title, subtitle):
    """新建演示文稿（python-pptx 生成，WPS 兼容格式）。"""
    uj = ctx.obj["json"]
    output = _abs(output)
    try:
        from pptx import Presentation
        from pptx.util import Pt
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        if subtitle and slide.placeholders[1]:
            slide.placeholders[1].text = subtitle
        prs.save(output)
        size = os.path.getsize(output)
        _out({"path": output, "title": title, "slides": 1, "size_bytes": size},
             f"✓ 已创建: {output}", uj)
    except Exception as e:
        click.echo(f"错误: {e}", err=True); sys.exit(1)


@impress.command("add-slide")
@click.argument("pptpath")
@click.option("-t", "--title", required=True, help="幻灯片标题")
@click.option("--body", default="", help="正文内容")
@click.pass_context
def impress_add_slide(ctx, pptpath, title, body):
    """追加幻灯片。"""
    uj = ctx.obj["json"]
    pptpath = _abs(pptpath)
    try:
        from pptx import Presentation
        prs = Presentation(pptpath)
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        if body and len(slide.placeholders) > 1:
            slide.placeholders[1].text = body
        prs.save(pptpath)
        _out({"path": pptpath, "title": title, "total_slides": len(prs.slides)},
             f"✓ 已添加幻灯片: {title}", uj)
    except Exception as e:
        click.echo(f"错误: {e}", err=True); sys.exit(1)


@impress.command("to-pdf")
@click.argument("pptpath")
@click.option("-o", "--output", default=None)
@click.pass_context
def impress_to_pdf(ctx, pptpath, output):
    """将演示文稿用 WPS wpp.exe 渲染为 PDF。"""
    uj = ctx.obj["json"]
    pptpath = _abs(pptpath)
    if not output:
        output = str(Path(pptpath).with_suffix(".pdf"))
    output = _abs(output)
    try:
        wpp = r"C:\Users\PRO\AppData\Local\Kingsoft\WPS Office\12.1.0.24034\office6\wpp.exe"
        if not os.path.exists(wpp):
            wpp = _find_wps_exe().replace("wps.exe", "wpp.exe")
        import subprocess
        # WPS 支持 /CONVERT-TO 参数
        result = subprocess.run(
            [wpp, pptpath, "/CONVERT-TO", "pdf", output],
            capture_output=True, timeout=30
        )
        # 若不支持命令行，用 COM 备用方案
        if not os.path.exists(output):
            import pywintypes, pythoncom, win32com.client as win32
            import winreg
            k = winreg.OpenKey(winreg.HKEY_CURRENT_USER,
                               r"SOFTWARE\Classes\KWPP.Application\CLSID")
            clsid = pywintypes.IID(winreg.QueryValueEx(k, "")[0])
            obj = pythoncom.CoCreateInstance(clsid, None,
                  pythoncom.CLSCTX_ALL, pythoncom.IID_IDispatch)
            app = win32.Dispatch(obj)
            prs = app.Presentations.Open(pptpath)
            prs.SaveAs(output, 32)
            prs.Close()
        size = os.path.getsize(output)
        with open(output, "rb") as f:
            valid = f.read(4) == b"%PDF"
        _out({"input": pptpath, "output": output,
              "size_bytes": size, "valid_pdf": valid},
             f"✓ PDF: {output} ({size:,} bytes)", uj)
    except Exception as e:
        click.echo(f"错误: {e}", err=True); sys.exit(1)


# ── 通用转换 ────────────────────────────────────────────────────
@cli.command("convert")
@click.argument("input_path")
@click.argument("output_path")
@click.pass_context
def convert(ctx, input_path, output_path):
    """智能格式转换（自动识别文档类型）。
    
    示例：wps-cli convert report.docx report.pdf
    """
    uj = ctx.obj["json"]
    ext = Path(input_path).suffix.lower()
    if ext in (".docx", ".doc", ".odt", ".rtf", ".txt"):
        ctx.invoke(writer_to_pdf, docpath=input_path, output=output_path)
    elif ext in (".xlsx", ".xls", ".ods", ".csv"):
        ctx.invoke(calc_to_pdf, xlspath=input_path, output=output_path)
    elif ext in (".pptx", ".ppt", ".odp"):
        ctx.invoke(impress_to_pdf, pptpath=input_path, output=output_path)
    else:
        click.echo(f"不支持的格式: {ext}", err=True); sys.exit(1)


# ── 版本信息 ─────────────────────────────────────────────────────
@cli.command("version")
@click.pass_context
def version(ctx):
    """显示 WPS 版本信息。"""
    uj = ctx.obj["json"]
    try:
        app = _get_com_app("writer")
        ver = app.Version
        build = getattr(app, "Build", "N/A")
        app.Quit()
        _out({"version": ver, "build": build, "com_id": "KWPS.Application"},
             f"WPS Version: {ver}", uj)
    except Exception as e:
        click.echo(f"错误: {e}", err=True)


if __name__ == "__main__":
    cli(obj={})
